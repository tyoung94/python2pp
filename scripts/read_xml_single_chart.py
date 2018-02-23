# -*- coding: utf-8 -*-
"""
Created on Mon Jan 22 10:59:04 2018

@author: d1tay01
"""
from pptx import Presentation
from pprint import pprint

prs = Presentation(r'C:\Users\D1TAY01\Desktop\pptx_desired_result_final.pptx')

slides = prs.slides

for slide in slides:
    shapes = slide.shapes
    for shape in shapes:
        x = shape.chart
        y = x.plots._plotArea
        with open(r'C:\Users\D1TAY01\Desktop\xml_desired_result_final', 'wt') as out:
            pprint((y.getparent()).getparent().xml, stream=out)