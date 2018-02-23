# -*- coding: utf-8 -*-
"""
Created on Thu Dec  7 13:32:27 2017

@author: d1tay01
"""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.util import Inches
import Haver
import pandas as pd
from copy import deepcopy
from pprint import pprint
from lxml import etree
from pptx.util import Pt
import os
import sys

os.chdir('H:/python2pp')
haver_path = r'O:/Haver/Data'
Haver.path(haver_path)

def clearSlide(slide):
    shapes = slide.shapes
    for shape in shapes:
        x = shape._element
        x.getparent().remove(x)
        
def buildChartData(series, database, start_date):
    df = Haver.data(series, database, start_date)
    ##Temp fix
    if series.__contains__('csent'):
        df = df[:-1]
    df = df.asfreq('D', how='start')
    df.index = df.index.strftime('%m/%d/%Y')
    df.index = pd.to_datetime(df.index)
    chart_data = CategoryChartData()
    chart_data.categories = df.index
    
    for x in series:
        chart_data.add_series(x, df[x])
    
    return df, chart_data

def insertLineChart(dim, chart_data):
    x = Inches(dim[0])
    y = Inches(dim[1])
    cx = Inches(dim[2])
    cy = Inches(dim[3])
    
    chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart
    return chart

def insertAreaChart(dim, chart_data):
    x = Inches(dim[0])
    y = Inches(dim[1])
    cx = Inches(dim[2])
    cy = Inches(dim[3])
    
    chart = slide.shapes.add_chart(
            XL_CHART_TYPE.AREA, x, y, cx, cy, chart_data
            ).chart
    return chart

def printTestXML(element):
    with open(r'./output/xml_test', 'wt') as out:
        pprint(element.xml, stream=out)
        

SLIDE_CHOICE = 12
start = '1995-01-01'
#These two prefixs are required for XML. I use whichever one the XML template specifies
prefix = '{http://schemas.openxmlformats.org/drawingml/2006/chart}'
prefixa = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

#Open BoD template
try:
    prs = Presentation(r'./templates/fomc_template.pptx')
except:
    sys.exit("Program terminating. PowerPoint template not found")
    
#set slide format
slide_layout = prs.slide_layouts[SLIDE_CHOICE]
#insert slide with selected format
slide = prs.slides.add_slide(slide_layout)
#When inserting a blank slide the Fed watermark is not included. I insert a nonblank slide and clear it as a work around
#clearSlide(slide)

#Build the consumer sentiment chart data
#Hard coded consumer sentiment indictors
series = ['csent', 'ccin']
cons_conf, chart_data = buildChartData(series, 'usecon', start)

#Dimensions for charts. I'd like to convert these to tuples, but i don't know how to efficiently unpack them and apply the "Inches" function found in the "insertLineChart" function
dim = [.5, .45, 8, 5]
main_chart = insertLineChart(dim, chart_data)

#build the recession chart data
recess = ['recessm']
rec, rec_chart_data = buildChartData(recess, 'usecon', start)
dim = [.5, .45, 4, 4]
rec_chart = insertAreaChart(dim, rec_chart_data)

#Set the scale of the y axis before copying
valAx = rec_chart.value_axis
valAx.maximum_scale = 0.1
valAx.minimum_scale = 0
valticks = valAx.tick_labels

#We don't want to see the secondary y axis values when we're doing recession shading
font = valticks.font
font.size = Pt(1)

#Format the recession series before copying
rec_series = rec_chart.series[0]
rec_fill = rec_series.format
fill = rec_fill.fill

########### This is the portal to magic XML fantasy land ##############
#Extract XML elements for the line chart and the area chart
rec_element = rec_chart.plots._plotArea
main_element = main_chart.plots._plotArea
############ Code changes gear after this point ##########################

# From the XML Final Template we know there are 4 "primary" elements missing in the main chart plot_area element
#We insert the elements in this order. This order was pulled from the desired XML result found in the "templates" folder
main_element.insert(0, rec_element.find(prefix + 'layout'))
main_element.insert(1, rec_element.find(prefix + 'areaChart'))
main_element.append(rec_element.find(prefix+'valAx'))
main_element.append(rec_element.find(prefix+'dateAx'))

#Now in the main_element SubElements, there are two valAx's and two dateAx's. I don't know how to use "find' when there are multiple subelements with the same name
#Because we appended the new axes to the end of the list, we can get them from the same place
children_elements = main_element.getchildren()
children_names = ['layout', 'rec_chart', 'main_chart', 'main_dateAx', 'main_valAx', 'rec_valAx', 'rec_dateAx']

main_dict = dict(zip(children_names, children_elements))

areaChart_copy = main_element.find(prefix + 'areaChart')
newDateAx = children_elements[-1]
newValAx = children_elements[-2]
oldValAx = children_elements[-3]
oldDateAx = children_elements[-4]
#################################################################################
#I don't know what these do. Just trying to match the desired xml


ser = areaChart_copy.find(prefix+'ser')
idx = ser.find(prefix+'idx')
idx.set('val', '2')
order = ser.find(prefix+'order')
order.set('val', '2')

#Format areachart fill and transparency
spPr = ser.find(prefix+'spPr')
solidFill = etree.SubElement(spPr, prefixa+'solidFill')
schemeClr = etree.SubElement(solidFill, prefixa+'schemeClr')
schemeClr.set('val', 'bg1')
lumMod = etree.SubElement(schemeClr, prefixa+'lumMod')
lumMod.set('val', '65000')
alpha = etree.SubElement(schemeClr, prefixa+'alpha')
alpha.set('val', '80000')

#format gridlines
grids = oldValAx.find(prefix+'majorGridlines')
spPr = etree.SubElement(grids, prefix + 'spPr')
ln = etree.SubElement(spPr, prefixa + 'ln')
dash = etree.SubElement(ln, prefixa + 'prstDash')
dash.set('val', 'lgDash')

valAx = main_element.find(prefix+'valAx')
valAx_scaling = valAx.find(prefix+'scaling')
orientation = etree.SubElement(valAx_scaling, prefix+'orientation')
orientation.set('val', 'minMax')

crossBetween = etree.SubElement(valAx, prefix+'crossBetween')
crossBetween.set('val', 'between')

numFormat = etree.SubElement(valAx, prefix+'numFmt')
numFormat.set('formatCode', 'General')
numFormat.set('sourceLinked', '1')

axPos = newValAx.find(prefix+'axPos')
axPos.set('val', 'r')
cross = newValAx.find(prefix+'crosses')
cross.set('val', 'max')
crossB = newValAx.find(prefix+'crossBetween')
crossB.set('val', 'between')
newValAx.remove(newValAx.find(prefix+'majorGridlines'))

delete = newDateAx.find(prefix+'delete')
delete.set('val', '1')
time = newDateAx.find(prefix+'baseTimeUnit')
time.set('val', 'months')
newDateAx.remove(newDateAx.find(prefix+'crosses'))

copy = deepcopy(numFormat)
numFormat.getparent().remove(numFormat)
valAx.insert(5, copy)

lineChart = main_element.find(prefix+'lineChart')
dLbls = etree.SubElement(lineChart, prefix+'dLbls')
legKey = etree.SubElement(dLbls, prefix+'showLegendKey')
legKey.set('val', '0')
showVal = etree.SubElement(dLbls, prefix+'showVal')
showVal.set('val', '0')
showCat = etree.SubElement(dLbls, prefix+'showCatName')
showCat.set('val', '0')
showSer = etree.SubElement(dLbls, prefix+'showSerName')
showSer.set('val', '0')
showPerc = etree.SubElement(dLbls, prefix+'showPercent')
showPerc.set('val', '0')
showBubs = etree.SubElement(dLbls, prefix+'showBubbleSize')
showBubs.set('val', '0')

dLbls_copy = deepcopy(dLbls)
dLbls.getparent().remove(dLbls)
lineChart.insert(4, dLbls_copy)

chartParent = main_element.getparent()
parent = (main_element.getparent()).getparent()
lang = etree.SubElement(parent, prefix+'lang')
lang.set('val', 'en-US')
rounded = etree.SubElement(parent, prefix+'roundedCorners')
rounded.set('val', '1')
style = etree.SubElement(parent, prefix+'style')
style.set('val', '2')

lang_copy = deepcopy(lang)
lang.getparent().remove(lang)
rounded_copy = deepcopy(rounded)
rounded.getparent().remove(rounded)
style_copy = deepcopy(style)
style.getparent().remove(style)

parent.insert(1, lang_copy)
parent.insert(2, rounded_copy)
parent.insert(3, style_copy)
#############################################################################

#Finally we delete the original recession chart. Because it is the last shape that was added to the chart, it'll be the last in the list of shapes on the slide
shapes = slide.shapes
x = shapes[-1]._element
x.getparent().remove(x)


#---------Format X axis------------------------
cat_axis = main_chart.category_axis
cat_axis.has_major_gridlines = False

cat_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW

tick_labels = cat_axis.tick_labels
tick_labels.number_format = 'mmmyy'
tick_labels.number_format_is_linked = False

##-------Remove Legend-------------------
main_chart.has_legend = False

##------Add title-----------------------
# =============================================================================
# w = h = Inches(1)
# t = l = Inches(0)
# title = slide.shapes.add_textbox(l,t,w,h)
# tf = title.text_frame
# tf.text = 'Consumer Sentiment'
# =============================================================================

printTestXML((main_element.getparent()).getparent())

prs.save(r'./output/fomc_template_mod.pptx')
 




##get Axis Ids of line chart
#lineValAxId = retrieveAxId(main_element, 'valAx')
#lineDateAxId = retrieveAxId(main_element, 'dateAx')


# =============================================================================
# def retrieveAxId(chart_element, axType):
#     ax = chart_element.find(prefix+axType)
#     axId = ax.find(prefix+'axId')
#     temp = axId.values()
#     val = temp[0]
#     return val
# 
# def correctIds(chart_element, valId, dateId):
#     children = areaChart_copy.getchildren()
#     areaChart_copy_dateAx = children[-2]
#     areaChart_copy_dateAx.set('val',dateId)
#     areaChart_copy_valAx = children[-1]
#     areaChart_copy_valAx.set('val', valId)
# 
# =============================================================================

















